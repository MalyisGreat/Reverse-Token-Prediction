from __future__ import annotations

import json
import math
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DATA_PATH = ROOT / "docs" / "results" / "reverse_runpod_8xh100_apr2026.json"
OUT_DIR = ROOT / "docs" / "results"
ASSET_DIR = OUT_DIR / "assets"

BG = "#07111f"
PANEL = "#0f1c33"
GRID = "#32405f"
TEXT = "#f5f7fb"
MUTED = "#9fb0ce"
ACCENT = "#67e8f9"
ACCENT_2 = "#ff8a5b"
ACCENT_3 = "#b794f6"
SUCCESS = "#34d399"


def load_data() -> dict:
    return json.loads(DATA_PATH.read_text(encoding="utf-8"))


def ensure_dirs() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def style_axes(ax) -> None:
    ax.set_facecolor(BG)
    for spine in ax.spines.values():
        spine.set_color(GRID)
    ax.tick_params(colors=MUTED, labelsize=10)
    ax.grid(True, color=GRID, linewidth=0.8, alpha=0.35)
    ax.title.set_color(TEXT)
    ax.xaxis.label.set_color(MUTED)
    ax.yaxis.label.set_color(MUTED)


def build_validation_chart(data: dict) -> Path:
    rows = data["validation_bpb"]
    steps = [row["step"] for row in rows]
    values = [row["bpb"] for row in rows]

    fig, ax = plt.subplots(figsize=(12, 7), dpi=200)
    fig.patch.set_facecolor(BG)
    style_axes(ax)
    ax.plot(steps, values, color=ACCENT, linewidth=3, marker="o", markersize=6)
    ax.fill_between(steps, values, max(values) + 0.05, color=ACCENT, alpha=0.08)
    ax.set_title("Reverse Validation BPB Across the Surviving Run", fontsize=20, fontweight="bold", pad=18)
    ax.set_xlabel("Training Step", fontsize=12)
    ax.set_ylabel("Validation Bits Per Byte", fontsize=12)
    ax.set_xlim(left=0)

    key_steps = {0, 1000, 1750, 3000, 4500}
    for step, value in zip(steps, values):
        if step in key_steps:
            ax.annotate(
                f"{step}\n{value:.3f}",
                (step, value),
                textcoords="offset points",
                xytext=(0, -24 if step == 0 else 10),
                ha="center",
                color=TEXT,
                fontsize=9,
                bbox=dict(boxstyle="round,pad=0.25", facecolor=PANEL, edgecolor=GRID, alpha=0.95),
            )

    ax.text(
        0.99,
        0.02,
        "Sources: surviving launch log + terminal transcript reconstruction",
        transform=ax.transAxes,
        ha="right",
        va="bottom",
        color=MUTED,
        fontsize=9,
    )
    out = ASSET_DIR / "validation_curve.png"
    fig.tight_layout()
    fig.savefig(out, facecolor=fig.get_facecolor(), bbox_inches="tight")
    plt.close(fig)
    return out


def build_improvement_chart(data: dict) -> Path:
    rows = data["validation_bpb"]
    labels = []
    drops = []
    for prev, curr in zip(rows, rows[1:]):
        labels.append(f"{prev['step']}→{curr['step']}")
        drops.append(prev["bpb"] - curr["bpb"])

    fig, ax = plt.subplots(figsize=(12, 6), dpi=200)
    fig.patch.set_facecolor(BG)
    style_axes(ax)
    colors = [SUCCESS if d >= 0 else ACCENT_2 for d in drops]
    bars = ax.bar(range(len(drops)), drops, color=colors, width=0.72)
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, rotation=35, ha="right", color=MUTED)
    ax.set_ylabel("BPB Drop Between Evaluations", fontsize=12)
    ax.set_title("Each Saved Validation Still Moved the Model", fontsize=20, fontweight="bold", pad=18)
    for bar, drop in zip(bars, drops):
        ax.text(
            bar.get_x() + bar.get_width() / 2,
            bar.get_height() + 0.005,
            f"{drop:.3f}",
            ha="center",
            va="bottom",
            color=TEXT,
            fontsize=9,
        )
    ax.axhline(0, color=GRID, linewidth=1.1)
    out = ASSET_DIR / "validation_improvements.png"
    fig.tight_layout()
    fig.savefig(out, facecolor=fig.get_facecolor(), bbox_inches="tight")
    plt.close(fig)
    return out


def build_probe_heatmap(data: dict) -> Path:
    rows = data["manual_probe_scores"]
    metrics = ["anchor_adherence", "topic_stability", "factuality", "repetition_control"]
    metric_labels = ["Anchor\nAdherence", "Topic\nStability", "Factuality", "Repetition\nControl"]
    matrix = [[row[m] for m in metrics] for row in rows]
    labels = [row["label"] for row in rows]

    fig, ax = plt.subplots(figsize=(12, 7), dpi=200)
    fig.patch.set_facecolor(BG)
    ax.set_facecolor(BG)
    im = ax.imshow(matrix, cmap="magma", vmin=1, vmax=5, aspect="auto")
    ax.set_xticks(range(len(metric_labels)))
    ax.set_xticklabels(metric_labels, color=MUTED, fontsize=10)
    ax.set_yticks(range(len(labels)))
    ax.set_yticklabels(labels, color=MUTED, fontsize=10)
    ax.set_title("Manual CPU Probe Scorecard", fontsize=20, fontweight="bold", color=TEXT, pad=16)
    for i, row in enumerate(matrix):
        for j, value in enumerate(row):
            ax.text(j, i, f"{value:.2f}".rstrip("0").rstrip("."), ha="center", va="center", color=TEXT, fontsize=9)
    cbar = fig.colorbar(im, ax=ax, fraction=0.04, pad=0.03)
    cbar.ax.yaxis.set_tick_params(color=MUTED)
    plt.setp(cbar.ax.get_yticklabels(), color=MUTED)
    cbar.set_label("Manual Score (1-5)", color=MUTED)
    ax.text(
        0.99,
        -0.13,
        "Based on checkpoint probes performed on CPU after the run lost remote GPU access.",
        transform=ax.transAxes,
        ha="right",
        va="top",
        color=MUTED,
        fontsize=9,
    )
    out = ASSET_DIR / "probe_tradeoffs.png"
    fig.tight_layout()
    fig.savefig(out, facecolor=fig.get_facecolor(), bbox_inches="tight")
    plt.close(fig)
    return out


def build_hero(data: dict) -> Path:
    width, height = 1600, 900
    image = Image.new("RGB", (width, height), BG)
    draw = ImageDraw.Draw(image)

    def font(size: int, bold: bool = False):
        candidates = [
            "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
            "C:/Windows/Fonts/calibrib.ttf" if bold else "C:/Windows/Fonts/calibri.ttf",
        ]
        for candidate in candidates:
            if Path(candidate).exists():
                return ImageFont.truetype(candidate, size=size)
        return ImageFont.load_default()

    title_font = font(66, bold=True)
    sub_font = font(30)
    stat_value_font = font(44, bold=True)
    stat_label_font = font(24)

    draw.text((92, 92), "Reverse Token Prediction", font=title_font, fill=TEXT)
    draw.text((96, 178), "RunPod 8xH100 reverse pretraining result report", font=sub_font, fill=MUTED)

    # Sparkline panel
    spark_x0, spark_y0, spark_w, spark_h = 880, 92, 620, 320
    draw.rounded_rectangle((spark_x0, spark_y0, spark_x0 + spark_w, spark_y0 + spark_h), radius=36, fill=PANEL, outline=GRID, width=2)
    steps = [row["step"] for row in data["validation_bpb"]]
    values = [row["bpb"] for row in data["validation_bpb"]]
    min_v, max_v = min(values), max(values)
    pts = []
    for step, value in zip(steps, values):
        x = spark_x0 + 44 + (spark_w - 88) * (step - min(steps)) / (max(steps) - min(steps))
        y = spark_y0 + spark_h - 44 - (spark_h - 88) * ((value - min_v) / (max_v - min_v))
        pts.append((x, y))
    for idx in range(len(pts) - 1):
        draw.line((pts[idx], pts[idx + 1]), fill=ACCENT, width=7)
    for x, y in pts:
        draw.ellipse((x - 7, y - 7, x + 7, y + 7), fill=ACCENT_2)
    draw.text((spark_x0 + 38, spark_y0 + 26), "Validation BPB", font=font(26, bold=True), fill=TEXT)
    draw.text((spark_x0 + 38, spark_y0 + 66), "3.168 → 0.757", font=font(44, bold=True), fill=ACCENT)

    stats = [
        ("1.384B", "parameters"),
        ("5.84B", "target train tokens"),
        ("0.757", "best preserved bpb"),
        ("0.97M", "steady-state tok/s"),
    ]
    card_w, card_h = 330, 156
    start_x = 92
    start_y = 322
    gap_x = 26
    gap_y = 26
    for idx, (value, label) in enumerate(stats):
        row = idx // 2
        col = idx % 2
        x0 = start_x + col * (card_w + gap_x)
        y0 = start_y + row * (card_h + gap_y)
        draw.rounded_rectangle((x0, y0, x0 + card_w, y0 + card_h), radius=28, fill=PANEL, outline=GRID, width=2)
        draw.text((x0 + 26, y0 + 28), value, font=stat_value_font, fill=TEXT)
        draw.text((x0 + 28, y0 + 96), label, font=stat_label_font, fill=MUTED)

    notes = [
        "Objective viability: confirmed",
        "Best decoding setting observed: 5000 default",
        "Main failure: checkpoints saved to ephemeral /root/.cache",
    ]
    notes_y = 684
    for idx, note in enumerate(notes):
        y = notes_y + idx * 46
        draw.rounded_rectangle((96, y + 4, 110, y + 18), radius=7, fill=ACCENT_2)
        draw.text((128, y - 6), note, font=font(28), fill=TEXT)

    out = ASSET_DIR / "github_hero.png"
    image.save(out)
    return out


def add_textbox(slide, left, top, width, height, text, font_size, color, bold=False, name="Arial", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
    return box


def add_bullets(slide, left, top, width, height, bullets, font_size=22, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
    return box


def set_slide_bg(slide, color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))


def build_presentation(data: dict, hero: Path, validation_chart: Path, improvement_chart: Path, heatmap: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    slide.shapes.add_picture(str(hero), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(6.5), Inches(0.8), "Experiment Setup", 28, TEXT, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.15), Inches(6.3), Inches(0.5), "What actually ran on the 8xH100 node", 15, MUTED)
    bullets = [
        "Vendored nanochat reverse run on 8x NVIDIA H100 80GB HBM3.",
        "1.384B-parameter causal transformer trained on BOS + reversed token rows.",
        "Target budget: 5.84B train tokens, 2,048 context, depth 24, 12 heads, 1,536 embedding width.",
        "Checkpoints saved every 1,000 steps, validation every 250 steps.",
        "Later transcript shows steady-state throughput near 0.97M tok/s and about 58-59% BF16 MFU."
    ]
    add_bullets(slide, Inches(0.75), Inches(1.75), Inches(5.6), Inches(3.8), bullets, font_size=20)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.05), Inches(5.3), Inches(4.9))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor.from_string(PANEL.lstrip("#"))
    panel.line.color.rgb = RGBColor.from_string(GRID.lstrip("#"))
    add_textbox(slide, Inches(7.35), Inches(1.4), Inches(4.4), Inches(0.5), "Preserved launch metrics", 22, TEXT, bold=True)
    details = [
        ("params", "1.384B"),
        ("target tokens", "5.84B"),
        ("early median tok/s", "560.7k"),
        ("later transcript band", "0.95M-0.98M"),
        ("best saved bpb", "0.757 at step 4500")
    ]
    for idx, (label, value) in enumerate(details):
        y = 2.0 + idx * 0.68
        add_textbox(slide, Inches(7.4), Inches(y), Inches(2.5), Inches(0.3), label.upper(), 10, MUTED, bold=True)
        add_textbox(slide, Inches(9.5), Inches(y - 0.05), Inches(2.1), Inches(0.4), value, 22, ACCENT, bold=True)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(6.5), Inches(0.8), "Validation Curve", 28, TEXT, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.05), Inches(6.5), Inches(0.5), "The run kept improving deep into the schedule.", 15, MUTED)
    slide.shapes.add_picture(str(validation_chart), Inches(0.7), Inches(1.45), width=Inches(7.0))
    slide.shapes.add_picture(str(improvement_chart), Inches(8.0), Inches(1.45), width=Inches(4.5))

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(6.8), Inches(0.8), "Generation Quality Trade-off", 28, TEXT, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.05), Inches(7.0), Inches(0.5), "CPU probes showed the objective worked, but truthfulness lagged behind anchor landing.", 15, MUTED)
    slide.shapes.add_picture(str(heatmap), Inches(0.7), Inches(1.45), width=Inches(7.0))
    probe_notes = [
        "3000 default: first checkpoint that looked genuinely useful.",
        "4000 default: better on some anchors, but repetition got worse.",
        "5000 default: best overall compromise.",
        "0.25 temperature: strongest loops.",
        "0.9 temperature: less looping, more hallucination.",
        "Bottom line: decoding needs repetition control, not just temperature tuning."
    ]
    add_bullets(slide, Inches(8.0), Inches(1.55), Inches(4.4), Inches(4.7), probe_notes, font_size=18)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(7.0), Inches(0.8), "What the Run Proved", 28, TEXT, bold=True)
    conclusion_bullets = [
        "Reverse-only pretraining is viable at nanochat scale: the model learned to generate plausible lead-ins for fixed suffix anchors.",
        "Validation BPB kept falling from 3.168 at step 0 to 0.757 at step 4500.",
        "The reverse objective improved topic steering and anchor adherence before it improved factual reliability.",
        "Hallucination was still expected at this parameter/data scale, so the right next move is better decoding or forward-model reranking, not abandoning the objective."
    ]
    add_bullets(slide, Inches(0.75), Inches(1.5), Inches(6.0), Inches(4.6), conclusion_bullets, font_size=21)
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.25), Inches(1.55), Inches(5.0), Inches(4.4))
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor.from_string(PANEL.lstrip("#"))
    callout.line.color.rgb = RGBColor.from_string(GRID.lstrip("#"))
    add_textbox(slide, Inches(7.55), Inches(1.9), Inches(4.1), Inches(0.4), "Best concise read", 18, ACCENT_2, bold=True)
    add_textbox(
        slide,
        Inches(7.55),
        Inches(2.45),
        Inches(4.1),
        Inches(2.6),
        "The run learned reverse landing behavior. It did not learn factual reconstruction. That is still a strong result, because it upgrades the idea from speculation to a working training objective.",
        24,
        TEXT,
    )

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(7.0), Inches(0.8), "Failure and Fix", 28, TEXT, bold=True)
    add_textbox(slide, Inches(0.7), Inches(1.05), Inches(7.5), Inches(0.5), "The remote run reached roughly 99.69% before the container restart wiped the weights.", 15, MUTED)
    left_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.55), Inches(5.65), Inches(4.8))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor.from_string(PANEL.lstrip("#"))
    left_panel.line.color.rgb = RGBColor.from_string(GRID.lstrip("#"))
    add_textbox(slide, Inches(1.05), Inches(1.9), Inches(4.7), Inches(0.4), "Root cause", 20, ACCENT_2, bold=True)
    add_bullets(
        slide,
        Inches(1.05),
        Inches(2.35),
        Inches(4.9),
        Inches(3.7),
        [
            "Checkpoints were written to /root/.cache/nanochat_reverse.",
            "The persistent volume and downloaded workspace only covered /workspace.",
            "When the container restarted, /root/.cache vanished and the trained weights were unrecoverable."
        ],
        font_size=19,
    )
    right_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.55), Inches(5.75), Inches(4.8))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RGBColor.from_string(PANEL.lstrip("#"))
    right_panel.line.color.rgb = RGBColor.from_string(GRID.lstrip("#"))
    add_textbox(slide, Inches(7.1), Inches(1.9), Inches(4.8), Inches(0.4), "Repo fix committed", 20, ACCENT, bold=True)
    add_bullets(
        slide,
        Inches(7.1),
        Inches(2.35),
        Inches(4.8),
        Inches(3.7),
        [
            "Reverse scripts now default to /workspace/nanochat_reverse when /workspace is writable.",
            "nanochat common utilities now prefer /workspace on Linux hosted environments before ~/.cache.",
            "Future RunPod launches will persist checkpoints by default instead of relying on ephemeral home storage."
        ],
        font_size=19,
    )

    out = OUT_DIR / "Reverse_Token_Prediction_Results_2026-04-29.pptx"
    prs.save(out)
    return out


def main() -> None:
    ensure_dirs()
    data = load_data()
    validation = build_validation_chart(data)
    improvements = build_improvement_chart(data)
    heatmap = build_probe_heatmap(data)
    hero = build_hero(data)
    pptx_path = build_presentation(data, hero, validation, improvements, heatmap)
    print(f"Wrote {validation}")
    print(f"Wrote {improvements}")
    print(f"Wrote {heatmap}")
    print(f"Wrote {hero}")
    print(f"Wrote {pptx_path}")


if __name__ == "__main__":
    main()
